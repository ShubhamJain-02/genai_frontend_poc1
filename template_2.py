from pptx import Presentation
import io
import pytest
import addphoto
# Load the existing presentation
def ppt(d,l,t):
    presentation = Presentation(t)
    def replace_text(shape, replacement_dict):
            if not shape.has_text_frame:
                return
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacement_dict.items():
                        if old_text == run.text:
                            run.text = run.text.replace(old_text, new_text)

        # Iterate through slides and shapes to replace text
    for slide in presentation.slides:
        for shape in slide.shapes:
            replace_text(shape, d)
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    for slide in range(l+1,13):
        xml_slides.remove(slides[slide])
    presentation.save('modified_presentation.pptx')
    
    return "modified_presentation.pptx"
